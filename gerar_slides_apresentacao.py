import os
import re
import ast

import pandas as pd
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

# --- 🎨 DESIGN SYSTEM V9 (LEGIBILIDADE + RESUMO HEURÍSTICO) ---
COR_VERDE_PRIMARIA = RGBColor(0, 102, 51)
COR_VERDE_ESCURO = RGBColor(0, 84, 42)
COR_TEXTO_HEADER = RGBColor(255, 255, 255)

COR_FUNDO_SLIDE = RGBColor(240, 242, 245)
COR_CARD_BG = RGBColor(255, 255, 255)
COR_BORDA_SUTIL = RGBColor(214, 220, 226)
COR_SOMBRA = RGBColor(120, 128, 140)

COR_TEXTO_CONTEUDO = RGBColor(48, 54, 61)
COR_TITULO_DARK = RGBColor(34, 40, 46)

FONTE_FAMILIA = "Segoe UI"
FONTE_TITULO_CARD = "Segoe UI Semibold"
FONTE_TITULO_PRINCIPAL = "Segoe UI Semibold"

ARQUIVO_ENTRADA = "data.csv"
ARQUIVO_SAIDA = "Apresentacao_Modern_V9.pptx"

COL_ACOES = "Actions"
COL_PROXIMAS = "Activities"
COL_PROBLEMAS = "Comments"

MAX_DESC_SENTENCAS = 5
MAX_DESC_CHARS = 800
MAX_CARD_CHARS = 900


PALAVRAS_CHAVE_PRIORITARIAS = {
    "objetivo",
    "status",
    "resultado",
    "impacto",
    "pendência",
    "pendencias",
    "bloqueio",
    "risco",
    "entrega",
    "prazo",
    "publicação",
    "integracao",
    "integração",
    "conclusão",
    "conclusao",
    "próximo",
    "proximo",
    "ação",
    "acao",
}


def limpar_html(texto):
    if pd.isna(texto) or str(texto).strip() == "":
        return " - "

    soup = BeautifulSoup(str(texto), "html.parser")
    texto_limpo = " ".join(soup.get_text(separator=" ").split())
    return texto_limpo


def carregar_dados_blindado(caminho_arquivo):
    if not os.path.exists(caminho_arquivo):
        print(f"ERRO: {caminho_arquivo} não encontrado.")
        return None

    encodings = ["utf-8", "cp1252", "latin1"]
    separadores = [",", ";"]

    for enc in encodings:
        for sep in separadores:
            try:
                pd.read_csv(caminho_arquivo, encoding=enc, sep=sep, nrows=2)
                print(f"Lido com sucesso: {enc} | separador '{sep}'")
                return pd.read_csv(caminho_arquivo, encoding=enc, sep=sep)
            except Exception:
                continue

    return None


def _normalizar_espacos(texto):
    return re.sub(r"\s+", " ", texto).strip()


def _normalizar_item_textual(item):
    """Normaliza item textual removendo artefatos comuns de serialização."""
    texto = _normalizar_espacos(str(item))
    return texto.strip(" \t\n\r\"'")


def _parse_lista_serializada(texto):
    """Converte strings no formato "['a', 'b']" em lista real, quando possível."""
    bruto = str(texto).strip()
    if not (bruto.startswith("[") and bruto.endswith("]")):
        return None

    try:
        valor = ast.literal_eval(bruto)
    except (ValueError, SyntaxError):
        return None

    if isinstance(valor, list):
        return [_normalizar_item_textual(v) for v in valor if _normalizar_item_textual(v)]

    return None


def _quebrar_sentencas(texto):
    partes = re.split(r"(?<=[\.!?;])\s+", texto)
    return [p.strip(" -\t\n") for p in partes if p and p.strip()]


def _pontuar_sentenca(sentenca):
    s_lower = sentenca.lower()
    score = 0

    for palavra in PALAVRAS_CHAVE_PRIORITARIAS:
        if palavra in s_lower:
            score += 4

    tamanho = len(sentenca)
    if 45 <= tamanho <= 180:
        score += 3
    elif 181 <= tamanho <= 260:
        score += 1
    elif tamanho < 25:
        score -= 2

    if any(char.isdigit() for char in sentenca):
        score += 1

    if ":" in sentenca:
        score += 1

    return score


def resumir_texto_heuristico(texto, max_sentencas=MAX_DESC_SENTENCAS, max_chars=MAX_DESC_CHARS):
    texto = _normalizar_espacos(str(texto))
    if not texto or texto == "-":
        return "-"

    sentencas = _quebrar_sentencas(texto)
    if not sentencas:
        return texto[:max_chars]

    if len(texto) <= max_chars and len(sentencas) <= max_sentencas:
        return texto

    sentencas_rank = []
    for idx, sentenca in enumerate(sentencas):
        score = _pontuar_sentenca(sentenca)
        # leve preferência por sentenças iniciais
        score += max(0, 2 - (idx // 3))
        sentencas_rank.append((score, idx, sentenca))

    selecionadas = sorted(sentencas_rank, key=lambda x: x[0], reverse=True)[:max_sentencas]
    selecionadas = sorted(selecionadas, key=lambda x: x[1])

    resumo = " ".join(s[2] for s in selecionadas)
    if len(resumo) > max_chars:
        resumo = resumo[:max_chars].rsplit(" ", 1)[0] + "..."

    return resumo


def transformar_em_bullets(texto, max_linhas=6):
    lista_serializada = _parse_lista_serializada(texto)
    if lista_serializada is not None:
        return lista_serializada[:max_linhas] if lista_serializada else ["-"]

    texto = _normalizar_espacos(str(texto))
    if not texto or texto == "-":
        return ["-"]

    if "\n" in str(texto):
        itens = [i.strip(" -\t") for i in str(texto).splitlines() if i.strip()]
        if itens:
            return itens[:max_linhas]

    sentencas = _quebrar_sentencas(texto)
    if not sentencas:
        return [_normalizar_item_textual(texto)]

    return [_normalizar_item_textual(s) for s in sentencas[:max_linhas] if _normalizar_item_textual(s)]


def _aplicar_texto_header(text_frame, titulo):
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = Cm(0.15)
    text_frame.margin_right = Cm(0.15)

    p_h = text_frame.paragraphs[0]
    p_h.text = str(titulo).upper()
    p_h.alignment = PP_ALIGN.CENTER
    p_h.space_before = Pt(0)
    p_h.space_after = Pt(0)

    p_h.font.name = FONTE_TITULO_CARD
    p_h.font.size = Pt(11.8)
    p_h.font.bold = True
    p_h.font.color.rgb = COR_TEXTO_HEADER


def _aplicar_texto_body(text_frame, conteudo, tamanho_fonte=11.5, bullets=False):
    text_frame.clear()
    text_frame.margin_left = Cm(0.42)
    text_frame.margin_right = Cm(0.42)
    text_frame.margin_top = Cm(0.26)
    text_frame.margin_bottom = Cm(0.24)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.word_wrap = True

    itens = conteudo if isinstance(conteudo, list) else [str(conteudo)]

    for i, item in enumerate(itens):
        if i == 0:
            p_b = text_frame.paragraphs[0]
        else:
            p_b = text_frame.add_paragraph()

        p_b.text = str(item)
        p_b.alignment = PP_ALIGN.LEFT
        p_b.space_after = Pt(2 if bullets else 1)
        p_b.space_before = Pt(0)
        p_b.line_spacing = 1.22
        p_b.level = 0

        if bullets and item.strip() not in {"-", ""}:
            p_b.text = f"• {item.strip()}"

        p_b.font.name = FONTE_FAMILIA
        p_b.font.size = Pt(tamanho_fonte)
        p_b.font.color.rgb = COR_TEXTO_CONTEUDO


def criar_card_moderno(
    slide,
    left,
    top,
    width,
    height,
    titulo,
    conteudo,
    is_html=False,
    tamanho_fonte=11.5,
    bullets=False,
):
    """Cria card com melhor legibilidade mantendo o layout original."""

    sombra = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left + Cm(0.07), top + Cm(0.07), width, height
    )
    sombra.fill.solid()
    sombra.fill.fore_color.rgb = COR_SOMBRA
    sombra.fill.transparency = 90
    sombra.line.fill.background()
    sombra.shadow.inherit = False
    sombra.adjustments[0] = 0.02

    h_header = Cm(0.9)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, h_header)
    header.fill.solid()
    header.fill.fore_color.rgb = COR_VERDE_PRIMARIA
    header.line.color.rgb = COR_VERDE_ESCURO
    header.line.width = Pt(0.5)
    _aplicar_texto_header(header.text_frame, titulo)

    h_body = height - h_header
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + h_header, width, h_body)
    body.fill.solid()
    body.fill.fore_color.rgb = COR_CARD_BG
    body.line.color.rgb = COR_BORDA_SUTIL
    body.line.width = Pt(0.9)

    if is_html:
        conteudo = limpar_html(conteudo)

    if isinstance(conteudo, list):
        conteudo_final = [_normalizar_item_textual(item) for item in conteudo if _normalizar_item_textual(item)]
        if not conteudo_final:
            conteudo_final = ["-"]
    else:
        conteudo_final = _normalizar_espacos(str(conteudo))
        if len(conteudo_final) > MAX_CARD_CHARS and not bullets:
            conteudo_final = conteudo_final[:MAX_CARD_CHARS].rsplit(" ", 1)[0] + "..."

    _aplicar_texto_body(
        body.text_frame,
        conteudo_final,
        tamanho_fonte=tamanho_fonte,
        bullets=bullets,
    )


def _adicionar_titulo_principal(slide, margin_x, width_total, txt_id, txt_titulo):
    shape_title = slide.shapes.add_textbox(margin_x, Cm(0.5), width_total, Cm(1.5))
    tf_title = shape_title.text_frame
    tf_title.clear()
    tf_title.margin_left = Cm(0.2)
    tf_title.margin_right = Cm(0.2)
    tf_title.vertical_anchor = MSO_ANCHOR.MIDDLE

    p_t = tf_title.paragraphs[0]
    p_t.space_before = Pt(0)
    p_t.space_after = Pt(0)

    run_id = p_t.add_run()
    run_id.text = f"ITEM {txt_id}"
    run_id.font.name = FONTE_TITULO_PRINCIPAL
    run_id.font.bold = True
    run_id.font.size = Pt(21)
    run_id.font.color.rgb = COR_VERDE_PRIMARIA

    run_sep = p_t.add_run()
    run_sep.text = "  |  "
    run_sep.font.size = Pt(20)
    run_sep.font.color.rgb = RGBColor(168, 176, 185)

    run_tit = p_t.add_run()
    run_tit.text = txt_titulo
    run_tit.font.name = "Segoe UI"
    run_tit.font.size = Pt(20)
    run_tit.font.color.rgb = COR_TITULO_DARK


def main():
    print("--- 🎨 GERADOR DE SLIDES V9 (LEGIBILIDADE + RESUMO) ---")

    df = carregar_dados_blindado(ARQUIVO_ENTRADA)
    if df is None:
        return

    for col in [COL_ACOES, COL_PROXIMAS, COL_PROBLEMAS]:
        if col not in df.columns:
            df[col] = ""
    df = df.fillna("")

    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    print(f"Processando {len(df)} itens...")

    for _, row in df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COR_FUNDO_SLIDE

        # --- GRID SYSTEM (inalterado) ---
        margin_x = Cm(1.0)
        width_total = Cm(31.867)
        gap = Cm(0.5)

        txt_id = str(row.get("ID", ""))
        txt_titulo = str(row.get("Title", ""))
        _adicionar_titulo_principal(slide, margin_x, width_total, txt_id, txt_titulo)

        current_y = Cm(2.2)

        # 2. Metadados
        h_meta = Cm(2.0)
        w_meta = (width_total - (3 * gap)) / 4
        meta_dados = [
            ("TIPO", row.get("Work Item Type", "")),
            ("RESPONSÁVEL", row.get("Assigned To", "").split("<")[0].strip()),
            ("STATUS", row.get("State", "")),
            ("PRIORIDADE", str(row.get("Priority", ""))),
        ]

        pos_x = margin_x
        for titulo, valor in meta_dados:
            criar_card_moderno(
                slide,
                pos_x,
                current_y,
                w_meta,
                h_meta,
                titulo,
                valor,
                tamanho_fonte=12.3,
            )
            pos_x += w_meta + gap

        current_y += h_meta + gap

        # 3. Descrição resumida heurística
        h_desc = Cm(4.5)
        descricao_limpa = limpar_html(row.get("Description", ""))
        descricao_resumida = resumir_texto_heuristico(descricao_limpa)
        bullets_desc = transformar_em_bullets(descricao_resumida, max_linhas=5)

        criar_card_moderno(
            slide,
            margin_x,
            current_y,
            width_total,
            h_desc,
            "DESCRIÇÃO DETALHADA",
            bullets_desc,
            tamanho_fonte=11.8,
            bullets=True,
        )

        current_y += h_desc + gap

        # 4. Ações e Próximas
        h_middle = Cm(5.0)
        w_half = (width_total - gap) / 2

        acoes_bullets = transformar_em_bullets(row[COL_ACOES], max_linhas=6)
        proximas_bullets = transformar_em_bullets(row[COL_PROXIMAS], max_linhas=6)

        criar_card_moderno(
            slide,
            margin_x,
            current_y,
            w_half,
            h_middle,
            "✅ AÇÕES REALIZADAS",
            acoes_bullets,
            tamanho_fonte=11.8,
            bullets=True,
        )
        criar_card_moderno(
            slide,
            margin_x + w_half + gap,
            current_y,
            w_half,
            h_middle,
            "📅 PRÓXIMAS ATIVIDADES",
            proximas_bullets,
            tamanho_fonte=11.8,
            bullets=True,
        )

        current_y += h_middle + gap

        # 5. Pontos de atenção
        h_total_slide = Cm(19.05)
        margin_bottom = Cm(1.0)
        h_restante = h_total_slide - current_y - margin_bottom
        if h_restante < Cm(2.0):
            h_restante = Cm(2.5)

        problemas_bullets = transformar_em_bullets(row[COL_PROBLEMAS], max_linhas=4)

        criar_card_moderno(
            slide,
            margin_x,
            current_y,
            width_total,
            h_restante,
            "⚠️ PONTOS DE ATENÇÃO",
            problemas_bullets,
            tamanho_fonte=11.8,
            bullets=True,
        )

    try:
        prs.save(ARQUIVO_SAIDA)
        print(f"\n✨ Apresentação salva com sucesso: {ARQUIVO_SAIDA}")
    except PermissionError:
        print(f"\n🚫 ERRO: O arquivo '{ARQUIVO_SAIDA}' está aberto. Feche-o e tente novamente.")


if __name__ == "__main__":
    main()
